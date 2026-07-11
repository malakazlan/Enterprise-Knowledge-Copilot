"""Tests for DOCX/PPTX/XLSX ingestion — real files through the real pipeline."""

from __future__ import annotations

import io
from collections.abc import Awaitable, Callable

import pytest
from httpx import AsyncClient

from app.models.user import User, UserRole

pytest.importorskip("docx")
pytest.importorskip("pptx")
pytest.importorskip("openpyxl")

MakeUser = Callable[..., Awaitable[User]]
AuthHeaders = Callable[..., Awaitable[dict[str, str]]]

DOCX_TYPE = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
XLSX_TYPE = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


def make_docx() -> bytes:
    from docx import Document

    document = Document()
    document.add_heading("Vacation Policy", 0)
    document.add_paragraph("Employees receive twenty five days of paid vacation per year.")
    table = document.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Tenure"
    table.rows[0].cells[1].text = "Extra days"
    table.rows[1].cells[0].text = "Five years"
    table.rows[1].cells[1].text = "Three"
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def make_pptx() -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Q3 Security Training"
    slide.placeholders[1].text = "Phishing reports doubled in the third quarter."
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def make_xlsx() -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "Budget"
    sheet.append(["Item", "Amount"])
    sheet.append(["Server hardware", 42000])
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


async def _admin(
    client: AsyncClient, make_user: MakeUser, auth_headers: AuthHeaders
) -> dict[str, str]:
    await make_user("admin@example.com", role=UserRole.ADMIN)
    return await auth_headers("admin@example.com")


async def _upload_and_ask(
    client: AsyncClient,
    headers: dict[str, str],
    filename: str,
    data: bytes,
    content_type: str,
    question: str,
    expected_fragment: str,
) -> None:
    upload = await client.post(
        "/api/v1/documents",
        headers=headers,
        files={"file": (filename, data, content_type)},
    )
    assert upload.status_code == 201, upload.text
    assert upload.json()["document"]["status"] == "completed"

    search = await client.post(
        "/api/v1/search", headers=headers, json={"query": question, "top_k": 3}
    )
    results = search.json()["results"]
    assert results, f"no retrieval results for {filename}"
    assert any(expected_fragment.lower() in r["content"].lower() for r in results)
    assert results[0]["filename"] == filename


async def test_docx_ingests_paragraphs_and_tables(
    client: AsyncClient, make_user: MakeUser, auth_headers: AuthHeaders
) -> None:
    headers = await _admin(client, make_user, auth_headers)
    await _upload_and_ask(
        client,
        headers,
        "policy.docx",
        make_docx(),
        DOCX_TYPE,
        "How many vacation days do employees get?",
        "twenty five days",
    )
    # Table content is searchable too.
    search = await client.post(
        "/api/v1/search", headers=headers, json={"query": "tenure extra days", "top_k": 3}
    )
    assert any("Five years | Three" in r["content"] for r in search.json()["results"])


async def test_pptx_pages_are_slides(
    client: AsyncClient, make_user: MakeUser, auth_headers: AuthHeaders
) -> None:
    headers = await _admin(client, make_user, auth_headers)
    await _upload_and_ask(
        client,
        headers,
        "training.pptx",
        make_pptx(),
        PPTX_TYPE,
        "What happened to phishing reports?",
        "doubled",
    )


async def test_xlsx_sheets_are_pages(
    client: AsyncClient, make_user: MakeUser, auth_headers: AuthHeaders
) -> None:
    headers = await _admin(client, make_user, auth_headers)
    await _upload_and_ask(
        client,
        headers,
        "budget.xlsx",
        make_xlsx(),
        XLSX_TYPE,
        "server hardware budget amount",
        "42000",
    )

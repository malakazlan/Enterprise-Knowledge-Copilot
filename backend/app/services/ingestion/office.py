"""Office document parsing: DOCX, PPTX, XLSX (the [office] extra).

Page mapping mirrors how people cite these formats: a PPTX page is a slide, an
XLSX page is a sheet, a DOCX is one logical page (Word files carry no page
boundaries in the file itself). Tables are flattened to pipe-separated rows so
their content is searchable and quotable.
"""

from __future__ import annotations

import asyncio
from pathlib import Path

from app.services.ingestion.ports import ParsedDocument, ParsedPage

_DOCX_TYPE = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_XLSX_TYPE = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"

_EXTENSIONS = {".docx": _DOCX_TYPE, ".pptx": _PPTX_TYPE, ".xlsx": _XLSX_TYPE}
_MAX_ROWS_PER_SHEET = 1000


def _row_text(cells: tuple[object, ...]) -> str:
    return " | ".join("" if cell is None else str(cell).strip() for cell in cells).strip(" |")


def _parse_docx(file_path: str) -> ParsedDocument:
    from docx import Document

    document = Document(file_path)
    lines: list[str] = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            text = _row_text(tuple(cell.text for cell in row.cells))
            if text:
                lines.append(text)
    title = (document.core_properties.title or "").strip() or (lines[0] if lines else None)
    return ParsedDocument(pages=[ParsedPage(page_number=1, text="\n".join(lines))], title=title)


def _parse_pptx(file_path: str) -> ParsedDocument:
    from pptx import Presentation

    presentation = Presentation(file_path)
    pages: list[ParsedPage] = []
    for number, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    lines.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    text = _row_text(tuple(cell.text for cell in row.cells))
                    if text:
                        lines.append(text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"Notes: {notes}")
        pages.append(ParsedPage(page_number=number, text="\n".join(lines)))
    return ParsedDocument(pages=pages, title=None)


def _parse_xlsx(file_path: str) -> ParsedDocument:
    from openpyxl import load_workbook

    workbook = load_workbook(file_path, read_only=True, data_only=True)
    try:
        pages: list[ParsedPage] = []
        for number, sheet in enumerate(workbook.worksheets, start=1):
            lines: list[str] = [f"Sheet: {sheet.title}"]
            for i, row in enumerate(sheet.iter_rows(values_only=True)):
                if i >= _MAX_ROWS_PER_SHEET:
                    lines.append("… (sheet truncated)")
                    break
                text = _row_text(row)
                if text:
                    lines.append(text)
            pages.append(ParsedPage(page_number=number, text="\n".join(lines)))
        return ParsedDocument(pages=pages, title=None)
    finally:
        workbook.close()


class OfficeParser:
    """DOCX/PPTX/XLSX parser; extraction runs in a worker thread (blocking IO)."""

    name = "office"

    def supports(self, content_type: str, filename: str) -> bool:
        if content_type in (_DOCX_TYPE, _PPTX_TYPE, _XLSX_TYPE):
            return True
        return Path(filename).suffix.lower() in _EXTENSIONS

    async def parse(self, *, file_path: str, content_type: str, filename: str) -> ParsedDocument:
        kind = _EXTENSIONS.get(Path(filename).suffix.lower(), content_type)
        if kind == _DOCX_TYPE:
            return await asyncio.to_thread(_parse_docx, file_path)
        if kind == _PPTX_TYPE:
            return await asyncio.to_thread(_parse_pptx, file_path)
        return await asyncio.to_thread(_parse_xlsx, file_path)
